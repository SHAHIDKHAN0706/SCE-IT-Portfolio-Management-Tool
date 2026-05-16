from __future__ import annotations

from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def build_executive_deck(items: list[dict]) -> bytes:
    prs = Presentation()
    sce_blue = RGBColor(0x00, 0x73, 0xCF)
    sce_navy = RGBColor(0x00, 0x2A, 0x3A)
    titles = [
        'Executive Summary', 'Portfolio Overview', 'Funding Gap Analysis', 'Value Stream Detail', 'Initiatives Register',
        'Decisions & Offsets', 'Cross-Portfolio Matrix', 'Risk & Dependencies', 'Priority Ranking — Top N',
        'Overlap & Consolidation Candidates', 'Recommendations', 'Appendix A', 'Appendix B'
    ]
    for title in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12), Inches(1.0))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = 'SCE IT Portfolio Management Tool'
        p.font.bold = True
        p.font.color.rgb = sce_blue
        p.font.size = Pt(16)
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.bold = True
        p2.font.color.rgb = sce_navy
        p2.font.size = Pt(26)
        b = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(10), Inches(0.6))
        b.text_frame.text = f'Records in scope: {len(items)}'
    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()
